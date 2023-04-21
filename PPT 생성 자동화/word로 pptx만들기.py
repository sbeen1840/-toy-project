
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document

#%%
from PIL import Image, ImageEnhance

# 이미지 파일 열기
image = Image.open('C:/Users/soop1/OneDrive/바탕 화면/dd.jpg')

# 밝기 보정을 위한 인스턴스 생성
enhancer = ImageEnhance.Brightness(image)

# 밝기 보정 적용 (0.0 ~ 1.0)
brightness = 0.4
image = enhancer.enhance(brightness)

# 이미지 파일 저장
image.save('dd_new.jpg')

#%%

# 파워포인트 프레젠테이션 생성
prs = Presentation()
prs.slide_width = 12192000
prs.slide_height = 6858000

# Microsoft Word 파일 열기
doc_path = 'C:/Users/soop1/OneDrive/바탕 화면/day2 .docx'
doc = Document(doc_path)

# Microsoft Word 파일에서 데이터를 읽어 파워포인트 슬라이드 생성
for para in doc.paragraphs:
    # 슬라이드 생성
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 삽입
    img_path = 'C:/Users/soop1/OneDrive/바탕 화면/dd_new.jpg'
    left = top = 0
    pic = slide.shapes.add_picture(img_path, left, top, prs.slide_width, prs.slide_height)

    # 텍스트 상자 생성
    textbox = slide.shapes.add_textbox(left=Cm(2), top=Cm(2), width=Cm(28), height=Cm(18))
    text_frame = textbox.text_frame

    # 텍스트 상자에 텍스트 추가
    text_frame.text = para.text.strip()

    # 텍스트 색상 변경
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = MSO_ANCHOR.MIDDLE  # 중앙 정렬
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE

        for run in paragraph.runs:
            font = run.font
            font.name = 'Calisto MT'
            font.size = Pt(60)
            font.color.rgb = RGBColor(255, 255, 255)  # 흰색

    # 텍스트 상자의 위치 설정
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.alignment = MSO_ANCHOR.MIDDLE
    text_frame.word_wrap = True

    # 텍스트 상자를 슬라이드 정중앙으로 이동
    textbox.left = int((prs.slide_width - textbox.width) / 2)
    textbox.top = int((prs.slide_height - textbox.height) / 2)

# 파워포인트 파일 저장
prs.save('presentation45ddddddddddd4.pptx')
